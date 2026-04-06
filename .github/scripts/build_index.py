"""Build the LlamaIndex vector index from study materials in the manifest.

Run after process_upload.py to rebuild the index with all current materials.
The index files are committed to the repo so the app can fetch them directly.

This builder uses normal PDF text extraction first, then falls back to OCR for
scanned pages and supplements low-text pages with sparse OCR labels. The sparse
OCR pass helps recover text from flowcharts, ER diagrams, and other visual pages
where the important signal is scattered around the page.
"""

import io
import json
import mimetypes
import os
import subprocess
import tempfile
import time
import urllib.error
import urllib.request
from pathlib import Path

import PyPDF2
import pypdfium2 as pdfium
import pytesseract
from llama_index.core import (
    Document,
    Settings,
    VectorStoreIndex,
    StorageContext,
    load_index_from_storage,
)
from llama_index.embeddings.huggingface import HuggingFaceEmbedding

MANIFEST_PATH = "manifest.json"
INDEX_DIR = "index"
INDEX_CACHE_PATH = "indexed_files.json"
EMBED_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
MIN_DIRECT_TEXT_CHARS = 120
OCR_RENDER_SCALE = 2
OCR_LANG = os.getenv("OCR_LANG", "eng")
OCR_TIMEOUT_SECONDS = int(os.getenv("OCR_TIMEOUT_SECONDS", "90"))
OFFICE_EXTENSIONS = {".doc", ".docx", ".ppt", ".pptx"}


def normalize_text(text):
    """Normalize text extracted from PDFs or OCR."""
    if not text:
        return ""
    filtered = text.encode("utf-16", "surrogatepass").decode("utf-16", "ignore")
    lines = [" ".join(line.split()) for line in filtered.splitlines()]
    return "\n".join(line for line in lines if line).strip()


def alpha_count(text):
    """Return the count of alpha-numeric characters in text."""
    return sum(1 for char in text if char.isalnum())


def read_pdf_content_safe(pdf_bytes):
    """Extract direct text from PDF bytes via PyPDF2."""
    pages = []
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(pdf_bytes))
        for page in reader.pages:
            try:
                page_text = normalize_text(page.extract_text())
                pages.append(page_text)
            except Exception as e:
                print(f"  Page extraction warning: {e}")
                pages.append("")
    except Exception as e:
        print(f"  PDF extraction error: {e}")
    return pages


def pil_image_from_pdf_page(page):
    """Render a PDF page to a PIL image for OCR."""
    bitmap = page.render(scale=OCR_RENDER_SCALE)
    return bitmap.to_pil()


def collect_sparse_ocr_lines(page_image):
    """Collect sparse OCR labels from visually dense pages."""
    data = pytesseract.image_to_data(
        page_image,
        lang=OCR_LANG,
        config="--oem 3 --psm 11",
        output_type=pytesseract.Output.DICT,
        timeout=OCR_TIMEOUT_SECONDS,
    )

    rows = {}
    count = len(data.get("text", []))
    for index in range(count):
        text = normalize_text(data["text"][index])
        confidence = data["conf"][index]
        if not text:
            continue
        try:
            conf_value = float(confidence)
        except (TypeError, ValueError):
            conf_value = -1
        if conf_value < 35:
            continue

        top = int(data["top"][index])
        line_key = round(top / 20)
        rows.setdefault(line_key, []).append((int(data["left"][index]), text))

    lines = []
    for _, words in sorted(rows.items()):
        line = " ".join(word for _, word in sorted(words))
        if line:
            lines.append(line)

    deduped = []
    seen = set()
    for line in lines:
        key = line.lower()
        if key in seen:
            continue
        seen.add(key)
        deduped.append(line)
    return deduped


def ocr_pdf_pages(pdf_bytes, direct_pages):
    """Run OCR on low-text pages and gather sparse labels for diagram pages."""
    combined_pages = []
    try:
        pdf = pdfium.PdfDocument(io.BytesIO(pdf_bytes))
    except Exception as e:
        print(f"  PDF render error: {e}")
        return "\n".join(page for page in direct_pages if page)

    try:
        for page_index in range(len(pdf)):
            direct_text = (
                direct_pages[page_index] if page_index < len(direct_pages) else ""
            )
            page = pdf[page_index]
            try:
                page_image = pil_image_from_pdf_page(page)
                direct_length = alpha_count(direct_text)

                parts = []
                if direct_text:
                    parts.append(direct_text)

                if direct_length < MIN_DIRECT_TEXT_CHARS:
                    ocr_text = normalize_text(
                        pytesseract.image_to_string(
                            page_image,
                            lang=OCR_LANG,
                            config="--oem 3 --psm 6",
                            timeout=OCR_TIMEOUT_SECONDS,
                        )
                    )
                    if ocr_text:
                        parts.append("[OCR TEXT]")
                        parts.append(ocr_text)

                sparse_lines = collect_sparse_ocr_lines(page_image)
                if sparse_lines:
                    sparse_block = "\n".join(f"- {line}" for line in sparse_lines[:50])
                    parts.append("[DIAGRAM LABELS]")
                    parts.append(sparse_block)

                page_text = "\n".join(part for part in parts if part).strip()
                if page_text:
                    combined_pages.append(f"[Page {page_index + 1}]\n{page_text}")
            except Exception as e:
                print(f"  OCR warning on page {page_index + 1}: {e}")
                if direct_text:
                    combined_pages.append(f"[Page {page_index + 1}]\n{direct_text}")
    finally:
        pdf.close()

    return "\n\n".join(combined_pages).strip()


def download_file(url, max_retries=3):
    """Download a file from URL with retry logic, handling HTTP 416 gracefully."""
    for attempt in range(max_retries):
        try:
            req = urllib.request.Request(url)
            with urllib.request.urlopen(req, timeout=60) as resp:
                return resp.read()
        except Exception as e:
            if isinstance(e, urllib.error.HTTPError) and e.code == 416:
                print(f"  Download got HTTP 416 (attempt {attempt + 1}/{max_retries}); retrying clean...")
                if attempt < max_retries - 1:
                    time.sleep(2 ** attempt)
                    continue
                raise
            print(f"  Download error (attempt {attempt + 1}): {e}")
            if attempt == max_retries - 1:
                raise
    return None


def get_mime_type(filename):
    """Guess MIME type from filename."""
    mime, _ = mimetypes.guess_type(filename)
    return mime or "application/octet-stream"


def get_extension(filename):
    """Return the lowercase file extension."""
    return Path(filename).suffix.lower()


def is_supported_file(filename, mime):
    """Return True when a file can be indexed."""
    extension = get_extension(filename)
    return (
        mime.startswith("text/")
        or mime == "application/pdf"
        or extension in OFFICE_EXTENSIONS
    )


def _pptx_to_pdf(file_bytes):
    """Convert PPTX bytes to PDF bytes using python-pptx + fpdf2."""
    from pptx import Presentation as PptxPresentation
    from pptx.util import Emu
    from fpdf import FPDF
    import io

    prs = PptxPresentation(io.BytesIO(file_bytes))

    slide_w_mm = prs.slide_width / Emu(914400) * 25.4
    slide_h_mm = prs.slide_height / Emu(914400) * 25.4

    pdf = FPDF(orientation="L", unit="mm", format=(slide_h_mm, slide_w_mm))
    pdf.set_auto_page_break(auto=True, margin=12)

    TITLE_SIZE = 18
    BODY_SIZE = 11
    MARGIN = 14

    for slide_idx, slide in enumerate(prs.slides, start=1):
        pdf.add_page()
        pdf.set_left_margin(MARGIN)
        pdf.set_right_margin(MARGIN)
        pdf.set_y(MARGIN)

        pdf.set_font("Helvetica", "I", 8)
        pdf.set_text_color(120, 120, 120)
        pdf.cell(0, 5, f"Slide {slide_idx}", ln=True)
        pdf.ln(2)

        title_text = ""
        body_parts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if not title_text and hasattr(shape, "placeholder_format"):
                        ph = shape.placeholder_format
                        if ph is not None and ph.idx in (0, 1):
                            title_text = text
                            continue
                    body_parts.append(text)

            if shape.shape_type == 13:
                try:
                    img_bytes = shape.image.blob
                    img_stream = io.BytesIO(img_bytes)
                    max_w = slide_w_mm - 2 * MARGIN
                    pdf.image(img_stream, x=MARGIN, w=min(max_w, 120))
                    pdf.ln(4)
                except Exception:
                    pass

        if title_text:
            pdf.set_font("Helvetica", "B", TITLE_SIZE)
            pdf.set_text_color(22, 49, 44)
            pdf.multi_cell(0, TITLE_SIZE * 0.5, title_text)
            pdf.ln(4)

        if body_parts:
            pdf.set_font("Helvetica", "", BODY_SIZE)
            pdf.set_text_color(50, 50, 50)
            for part in body_parts:
                pdf.multi_cell(0, BODY_SIZE * 0.45, part)
                pdf.ln(2)

    if len(prs.slides) == 0:
        pdf.add_page()
        pdf.set_font("Helvetica", "I", 12)
        pdf.cell(0, 10, "This presentation has no slides.", ln=True)

    return bytes(pdf.output())


def _docx_to_pdf(file_bytes):
    """Convert DOCX bytes to PDF bytes using python-docx + fpdf2."""
    from docx import Document as DocxDocument
    from fpdf import FPDF
    import io

    doc = DocxDocument(io.BytesIO(file_bytes))

    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    MARGIN = 16
    pdf.set_left_margin(MARGIN)
    pdf.set_right_margin(MARGIN)

    HEADING_SIZES = {"Heading 1": 20, "Heading 2": 16, "Heading 3": 14}
    BODY_SIZE = 11

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            pdf.ln(3)
            continue

        style_name = para.style.name if para.style else ""

        if style_name in HEADING_SIZES:
            size = HEADING_SIZES[style_name]
            pdf.set_font("Helvetica", "B", size)
            pdf.set_text_color(22, 49, 44)
            pdf.ln(4)
            pdf.multi_cell(0, size * 0.5, text)
            pdf.ln(3)
        else:
            is_bold = any(run.bold for run in para.runs if run.bold is not None)
            pdf.set_font("Helvetica", "B" if is_bold else "", BODY_SIZE)
            pdf.set_text_color(50, 50, 50)
            pdf.multi_cell(0, BODY_SIZE * 0.45, text)
            pdf.ln(1.5)

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_stream = io.BytesIO(rel.target_part.blob)
                pdf.image(img_stream, x=MARGIN, w=100)
                pdf.ln(4)
            except Exception:
                pass

    return bytes(pdf.output())


def convert_office_to_pdf(file_bytes, filename):
    """Convert an Office document to PDF bytes using pure Python."""
    ext = get_extension(filename)
    try:
        if ext in (".pptx", ".ppt"):
            return _pptx_to_pdf(file_bytes)
        elif ext in (".docx", ".doc"):
            return _docx_to_pdf(file_bytes)
    except Exception as err:
        print(f"  Office-to-PDF conversion error ({filename}): {err}")
    return None


def main():
    if not os.path.exists(MANIFEST_PATH):
        print("No manifest.json found. Nothing to index.")
        return

    with open(MANIFEST_PATH, "r") as f:
        manifest = json.load(f)

    if not manifest:
        print("Manifest is empty. Nothing to index.")
        return

    indexed_urls = []
    if os.path.exists(INDEX_CACHE_PATH):
        try:
            with open(INDEX_CACHE_PATH, "r") as f:
                indexed_urls = json.load(f)
        except Exception as e:
            print(f"Warning: could not read cache: {e}")

    # Prune stale cache entries whose URLs are no longer in the manifest.
    # This prevents dangling references when files are removed from the manifest.
    manifest_urls = {
        entry["download_url"]
        for subjects in manifest.values()
        for types in subjects.values()
        for file_list in types.values()
        for entry in file_list
    }
    original_cache_count = len(indexed_urls)
    indexed_urls = [u for u in indexed_urls if u in manifest_urls]
    pruned_count = original_cache_count - len(indexed_urls)
    if pruned_count > 0:
        print(f"Pruned {pruned_count} stale cache entries not in manifest.")

    # Collect all files from manifest
    documents = []
    newly_indexed_urls = []
    for semester, subjects in manifest.items():
        for subject, types in subjects.items():
            for file_type, file_list in types.items():
                for file_entry in file_list:
                    name = file_entry["name"]
                    url = file_entry["download_url"]

                    if url in indexed_urls:
                        print(f"Skipping already indexed file: {name}")
                        continue

                    mime = get_mime_type(name)

                    if not is_supported_file(name, mime):
                        print(f"Skipping unsupported: {name} ({mime})")
                        continue

                    print(f"Processing: {name}")
                    try:
                        prev_doc_count = len(documents)
                        content = download_file(url)
                        if not content:
                            continue

                        metadata = {
                            "filename": name,
                            "semester": semester,
                            "subject": subject,
                            "type": file_type,
                        }

                        if mime.startswith("text/"):
                            text = content.decode("utf-8", errors="ignore")
                            if text.strip():
                                documents.append(Document(text=text, metadata=metadata))

                        elif mime == "application/pdf":
                            direct_pages = read_pdf_content_safe(content)
                            text = ocr_pdf_pages(content, direct_pages)
                            if text.strip():
                                documents.append(Document(text=text, metadata=metadata))
                            else:
                                print(f"  No text extracted from: {name}")

                        elif get_extension(name) in OFFICE_EXTENSIONS:
                            pdf_bytes = convert_office_to_pdf(content, name)
                            if not pdf_bytes:
                                print(f"  Could not convert Office file: {name}")
                                continue

                            direct_pages = read_pdf_content_safe(pdf_bytes)
                            text = ocr_pdf_pages(pdf_bytes, direct_pages)
                            if text.strip():
                                documents.append(Document(text=text, metadata=metadata))
                            else:
                                print(f"  No text extracted after conversion: {name}")

                        if len(documents) > prev_doc_count:
                            newly_indexed_urls.append(url)

                    except Exception as e:
                        print(f"  Error processing {name}: {e}")

    if not documents:
        print("No new documents to index.")
        # Still persist the pruned cache so stale entries are removed on disk.
        if pruned_count > 0:
            with open(INDEX_CACHE_PATH, "w") as f:
                json.dump(indexed_urls, f, indent=2)
            print(f"Cache saved with {pruned_count} stale entries removed.")
        return

    print(f"\nBuilding/Updating index with {len(documents)} new documents...")
    embed_model = HuggingFaceEmbedding(model_name=EMBED_MODEL)
    Settings.embed_model = embed_model

    if os.path.exists(INDEX_DIR) and len(os.listdir(INDEX_DIR)) > 0:
        print("Loading existing index from storage...")
        storage_context = StorageContext.from_defaults(persist_dir=INDEX_DIR)
        index = load_index_from_storage(storage_context, embed_model=embed_model)
        for doc in documents:
            index.insert(doc)
    else:
        print("Creating new index...")
        index = VectorStoreIndex.from_documents(documents, embed_model=embed_model)

    index.storage_context.persist(persist_dir=INDEX_DIR)
    print(f"Index persisted to {INDEX_DIR}/")

    all_indexed = list(set(indexed_urls + newly_indexed_urls))
    with open(INDEX_CACHE_PATH, "w") as f:
        json.dump(all_indexed, f, indent=2)
    print(f"Cache updated with {len(newly_indexed_urls)} new URLs.")


if __name__ == "__main__":
    main()
